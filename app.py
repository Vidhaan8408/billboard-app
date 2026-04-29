
from flask import Flask, render_template, request, send_file, jsonify, redirect, url_for
import json
import os
import re
import threading
import uuid
from datetime import datetime
from zipfile import ZipFile
from concurrent.futures import ThreadPoolExecutor

from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation

app = Flask(__name__)

jobs = {}


def normalize_site_id(value):
    try:
        return str(int(float(value)))
    except (TypeError, ValueError):
        return re.sub(r"\.0$", "", str(value).strip())


def pick(s, *keys):
    for k in keys:
        if k in s and s[k] not in (None, ""):
            return s[k]
    return ""


def add_full_slide(prs, path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(path, 0, 0, width=prs.slide_width, height=prs.slide_height)


def fetch_image_bytes(url):
    try:
        from io import BytesIO
        from urllib.request import urlopen
        with urlopen(url, timeout=8) as response:
            return BytesIO(response.read())
    except Exception:
        return None


def natural_image_sort_key(filename):
    match = re.search(r"_(\d+)\.(jpg|jpeg|png|webp)$", filename, re.IGNORECASE)
    if match:
        return int(match.group(1))
    return 0


def form_get(form_data, key, default=""):
    values = form_data.get(key)
    if values is None:
        return default
    if isinstance(values, list):
        return values[0] if values else default
    return values


def form_getlist(form_data, key):
    values = form_data.get(key)
    if values is None:
        return []
    if isinstance(values, list):
        return values
    return [values]


def render_loading_page(job_id):
    return f"""
    <html>
    <head>
        <title>Generating</title>
        <meta name="viewport" content="width=device-width, initial-scale=1">
        <style>
        body {{
            font-family: Arial, sans-serif;
            text-align: center;
            margin-top: 40vh;
            padding: 20px;
        }}
        </style>
    </head>
    <body>
        <h2>Generating files...</h2>

        <script>
        async function poll() {{
            try {{
                const res = await fetch("/job-status/{job_id}", {{ cache: "no-store" }});
                const data = await res.json();

                if (data.status === "done") {{
                    window.location.href = "/result/{job_id}";
                    return;
                }}

                if (data.status === "error") {{
                    document.body.innerHTML = "<h2>Generation failed</h2><pre>" + (data.error || "Unknown error") + "</pre>";
                    return;
                }}
            }} catch (e) {{
                console.log(e);
            }}

            setTimeout(poll, 1000);
        }}

        poll();
        </script>
    </body>
    </html>
    """


def render_result_page(client_name):
    return f"""
    <html>
    <head>
    <title>Done</title>
    <meta name="viewport" content="width=device-width, initial-scale=1">

    <style>
    body {{
        font-family: Arial;
        text-align: center;
        margin-top: 40px;
        padding: 20px;
    }}

    .btn {{
        display: block;
        width: 80%;
        margin: 10px auto;
        padding: 12px;
        background: black;
        color: white;
        text-decoration: none;
        border-radius: 8px;
        border: none;
        font-size: 16px;
    }}

    .done-btn {{
        background: #28a745;
    }}
    </style>
    </head>

    <body>

    <h2>✅ Files Generated</h2>

    <a class="btn" href="/download/{client_name}?type=ppt">Download PPT</a>
    <a class="btn" href="/download/{client_name}?type=excel">Download Excel</a>

    <button class="btn" onclick="shareFiles()">Share</button>

    <button class="btn done-btn" onclick="goHome()">Done</button>

    <script>
    async function shareFiles() {{
        const pptUrl = window.location.origin + "/download/{client_name}?type=ppt";
        const excelUrl = window.location.origin + "/download/{client_name}?type=excel";
        const text = "PPT: " + pptUrl + "\\nExcel: " + excelUrl;

        try {{
            const pptRes = await fetch(pptUrl);
            const excelRes = await fetch(excelUrl);

            const pptBlob = await pptRes.blob();
            const excelBlob = await excelRes.blob();

            const pptFile = new File([pptBlob], "{client_name}.pptx", {{
                type: pptBlob.type || "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            }});
            const excelFile = new File([excelBlob], "{client_name}.xlsx", {{
                type: excelBlob.type || "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            }});

            if (navigator.canShare && navigator.canShare({{ files: [pptFile, excelFile] }})) {{
                await navigator.share({{
                    title: "Veena Advertising Files",
                    text: text,
                    files: [pptFile, excelFile]
                }});
                return;
            }}
        }} catch (err) {{
            console.log("File share failed, falling back:", err);
        }}

        try {{
            if (navigator.share) {{
                await navigator.share({{
                    title: "Veena Advertising Files",
                    text: text
                }});
                return;
            }}
        }} catch (err) {{
            console.log("Link share failed, falling back to clipboard:", err);
        }}

        try {{
            await navigator.clipboard.writeText(text);
            alert("Share not available here. The links were copied.");
        }} catch (err) {{
            alert(text);
        }}
    }}

    function goHome() {{
        window.location.replace("/");
    }}
    </script>

    </body>
    </html>
    """


with open("data.json", "r", encoding="utf-8") as f:
    sites = json.load(f)


@app.route("/")
def index():
    return render_template("index.html", sites=sites)


@app.route("/job-status/<job_id>")
def job_status(job_id):
    job = jobs.get(job_id)
    if not job:
        return jsonify({"status": "error", "error": "Job not found"}), 404

    return jsonify({
        "status": job.get("status", "running"),
        "error": job.get("error", "")
    })


@app.route("/loading/<job_id>")
def loading(job_id):
    job = jobs.get(job_id)
    if not job:
        return "Error: Job not found", 404

    return render_loading_page(job_id)


def process_generation(job_id, form_data, selected, mode, client_name, batch_date, batch_ids):
    try:
        os.makedirs("outputs", exist_ok=True)

        # ------------------------
        # EXCEL (FINAL FIXED FULL)
        # ------------------------
        wb = Workbook()
        ws = wb.active
        ws.title = "Sites"

        title_fill = PatternFill("solid", fgColor="FFD966")
        header_fill = PatternFill("solid", fgColor="BDD7EE")
        data_fill = PatternFill("solid", fgColor="FFF2CC")
        note_fill = PatternFill("solid", fgColor="FFF200")

        border = Border(
            left=Side(style="thin"),
            right=Side(style="thin"),
            top=Side(style="thin"),
            bottom=Side(style="thin")
        )

        extra_cols = 0
        if mode in ["rates", "both"]:
            extra_cols += 1
        if mode in ["availability", "both"]:
            extra_cols += 1

        max_cols = 8 + extra_cols

        # TITLE
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max_cols)
        ws["A1"] = "VEENA ADVERTISING"
        ws["A1"].font = Font(bold=True, size=16)
        ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
        ws["A1"].fill = title_fill
        ws["A1"].border = border
        ws.row_dimensions[1].height = 24
        ws.row_dimensions[2].height = 8

        # HEADERS
        ws["A3"] = "Sr. No."
        ws["B3"] = "Center"
        ws["C3"] = "Location"
        ws["D3"] = "Size"
        ws["D4"] = "W"
        ws["E4"] = "H"
        ws["F3"] = "Type"
        ws["G3"] = "Sq. Ft."
        ws["H3"] = "Lat. & Long."

        for row in [3, 4]:
            for col in range(1, 9):
                c = ws.cell(row=row, column=col)
                c.fill = header_fill
                c.border = border
                c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                c.font = Font(bold=True)

        ws.merge_cells("A3:A4")
        ws.merge_cells("B3:B4")
        ws.merge_cells("C3:C4")
        ws.merge_cells("D3:E3")
        ws.merge_cells("F3:F4")
        ws.merge_cells("G3:G4")
        ws.merge_cells("H3:H4")

        next_col = 9
        if mode in ["rates", "both"]:
            ws.cell(row=3, column=next_col, value="Rate per month")
            ws.merge_cells(start_row=3, start_column=next_col, end_row=4, end_column=next_col)
            next_col += 1

        if mode in ["availability", "both"]:
            ws.cell(row=3, column=next_col, value="Availability")
            ws.merge_cells(start_row=3, start_column=next_col, end_row=4, end_column=next_col)

        data_start_row = 6

        for idx, s in enumerate(selected, start=1):
            site_id = normalize_site_id(s.get("id"))

            center = pick(s, "center")
            location = pick(s, "location")
            w = pick(s, "W", "width")
            h = pick(s, "H", "height")
            type_ = pick(s, "type")
            sq_ft = pick(s, "sq_ft")
            lat_long = pick(s, "lat_long")
            rate = pick(s, "rate")

            row_values = [idx, center, location, w, h, type_, sq_ft, lat_long]

            if mode in ["rates", "both"]:
                row_values.append(rate)

            if mode in ["availability", "both"]:
                availability_choice = form_get(form_data, f"availability_{site_id}", "")
                date_val = form_get(form_data, f"date_{site_id}", "")

                if availability_choice == "hold":
                    availability = "On Hold"
                elif site_id in batch_ids and batch_date:
                    try:
                        availability = datetime.strptime(batch_date, "%Y-%m-%d").strftime("%d %B %Y")
                    except:
                        availability = ""
                elif availability_choice == "now":
                    availability = "Available Now"
                elif availability_choice == "date":
                    if date_val:
                        try:
                            availability = datetime.strptime(date_val, "%Y-%m-%d").strftime("%d %B %Y")
                        except:
                            availability = ""
                    else:
                        availability = ""
                else:
                    availability = ""

                row_values.append(availability)

            excel_row = data_start_row + idx - 1

            for col_idx, value in enumerate(row_values, start=1):
                c = ws.cell(row=excel_row, column=col_idx, value=value)
                c.fill = data_fill
                c.border = border
                c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

            ws.row_dimensions[excel_row].height = 18

        widths = {
            "A": 8,
            "B": 12,
            "C": 45,
            "D": 7,
            "E": 7,
            "F": 8,
            "G": 10,
            "H": 20,
            "I": 14,
            "J": 16
        }

        for col, width in widths.items():
            if ord(col) - 64 <= max_cols:
                ws.column_dimensions[col].width = width

        note_row = data_start_row + len(selected) + 2

        ws.cell(note_row, 2, "Notes:-")
        ws.cell(note_row, 2).fill = note_fill
        ws.cell(note_row, 2).font = Font(bold=True, color="FF0000")
        ws.cell(note_row, 2).alignment = Alignment(horizontal="center")
        ws.cell(note_row, 2).border = border

        notes = [
            "All the sites are subject to availability at the time of confirmation.",
            "Mail Confirmation or Purchase Order is mandatory to start the campaign.",
            "Flex printing mounting will be charged extra.",
            "Payment will be required in advance.",
            "GST (18%) will be charged.",
            "Please note that available date may change in case of present display renewal.",
            "For any reasons, if your flex is damaged, it will be your responsibility to provide us with a new flex.",
            "Please Check Availability At The Final Confirmation Time",
        ]

        for i, text in enumerate(notes, start=1):
            r = note_row + i
            ws.cell(r, 2, "•")
            ws.cell(r, 2).alignment = Alignment(horizontal="center", vertical="center")
            ws.cell(r, 2).border = border

            ws.merge_cells(start_row=r, start_column=3, end_row=r, end_column=max_cols)
            ws.cell(r, 3, text)
            ws.cell(r, 3).alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)

        footer = [
            "Thanking you,",
            "Veena Advertising",
            "Shishir Chauhan",
            "9824234436",
            "veenaadvertising@gmail.com"
        ]

        footer_start = note_row + len(notes) + 2
        for i, line in enumerate(footer):
            r = footer_start + i
            ws.merge_cells(start_row=r, start_column=3, end_row=r, end_column=max_cols)
            c = ws.cell(r, 3, line)
            c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
            if line == "veenaadvertising@gmail.com":
                c.font = Font(color="0000FF", underline="single")

        excel_path = f"outputs/{client_name}.xlsx"
        wb.save(excel_path)

        # ------------------------
        # PPT
        # ------------------------
        prs = Presentation()

        from io import BytesIO
        from urllib.request import urlopen

        cloud_name = os.environ.get("CLOUDINARY_CLOUD_NAME", "dub8ndson")

        # INTRO SLIDE
        intro_url = f"https://res.cloudinary.com/{cloud_name}/image/upload/intro.jpg"
        try:
            with urlopen(intro_url, timeout=20) as response:
                add_full_slide(prs, BytesIO(response.read()))
        except:
            pass

        # MAIN SLIDES
        for s in selected:
            sid = normalize_site_id(s.get("id"))

            for idx in range(1, 6):
                img = None

                for ext in ("jpg", "jpeg", "png"):
                    url = f"https://res.cloudinary.com/{cloud_name}/image/upload/{sid}_{idx}.{ext}"
                    img = fetch_image_bytes(url)

                    if img:
                        break   # STOP after first valid image

                if img:
                    add_full_slide(prs, img)


        # THANK YOU SLIDE
        thankyou_url = f"https://res.cloudinary.com/{cloud_name}/image/upload/thankyou.jpg"
        try:
            with urlopen(thankyou_url, timeout=20) as response:
                add_full_slide(prs, BytesIO(response.read()))
        except:
            pass

        ppt_path = f"outputs/{client_name}.pptx"
        prs.save(ppt_path)

        zip_path = f"outputs/{client_name}.zip"
        with ZipFile(zip_path, "w") as zipf:
            zipf.write(excel_path, os.path.basename(excel_path))
            zipf.write(ppt_path, os.path.basename(ppt_path))

        jobs[job_id]["status"] = "done"
        jobs[job_id]["error"] = ""

    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        print(f"Error in job {job_id}: {e}")


@app.route("/generate", methods=["POST"])
def generate():
    form_data = request.form.to_dict(flat=False)

    client_name = form_get(form_data, "client_name", "").strip()
    client_name = re.sub(r"[^a-zA-Z0-9_\- ]", "", client_name).replace(" ", "_")

    if not client_name or client_name in [".", ".."]:
        return "Error: Invalid client name"

    mode = form_get(form_data, "mode", "standard")

    selected_ids = set()

    # STANDARD + RATES
    if mode in ["standard", "rates"]:
        selected_ids = set(form_getlist(form_data, "site"))

    # AVAILABILITY + BOTH
    else:
        for s in sites:
            site_id = normalize_site_id(s.get("id"))
            availability_choice = form_get(form_data, f"availability_{site_id}", "")

            if availability_choice in ["now", "date", "hold"]:
                selected_ids.add(site_id)

        selected_ids |= set(form_getlist(form_data, "batch_site"))

    if not selected_ids:
        return "Error: No sites selected"

    selected = [
        s for s in sites
        if normalize_site_id(s.get("id")) in selected_ids
    ]

    if not selected:
        return "Error: No valid sites selected"

    batch_date = form_get(form_data, "batch_date", "").strip()
    batch_ids = set(form_getlist(form_data, "batch_sites")) | set(form_getlist(form_data, "batch_site"))

    job_id = uuid.uuid4().hex
    jobs[job_id] = {
        "status": "running",
        "client_name": client_name,
        "error": ""
    }

    thread = threading.Thread(
        target=process_generation,
        args=(job_id, form_data, selected, mode, client_name, batch_date, batch_ids),
        daemon=True
    )
    thread.start()

    return redirect(url_for("loading", job_id=job_id))


@app.route("/result/<job_id>")
def result(job_id):
    job = jobs.get(job_id)
    if not job:
        return "Error: Job not found", 404

    if job.get("status") == "error":
        return f"""
        <html>
        <head>
            <title>Error</title>
            <meta name="viewport" content="width=device-width, initial-scale=1">
        </head>
        <body style="font-family: Arial; text-align:center; margin-top:40px;">
            <h2>Generation failed</h2>
            <p>{job.get("error", "Unknown error")}</p>
            <a href="/" style="display:inline-block; margin-top:20px;">Go Home</a>
        </body>
        </html>
        """

    if job.get("status") != "done":
        return redirect(url_for("loading", job_id=job_id))

    return render_result_page(job["client_name"])


@app.route("/download/<client_name>")
def download(client_name):
    file_type = request.args.get("type")

    if file_type == "ppt":
        path = f"outputs/{client_name}.pptx"
        return send_file(
            path,
            as_attachment=True,
            download_name=f"{client_name}.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    elif file_type == "excel":
        path = f"outputs/{client_name}.xlsx"
        return send_file(
            path,
            as_attachment=True,
            download_name=f"{client_name}.xlsx",
            mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    else:
        path = f"outputs/{client_name}.zip"
        return send_file(
            path,
            as_attachment=True,
            download_name=f"{client_name}.zip",
            mimetype="application/zip"
        )


if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=5050, use_reloader=False)
